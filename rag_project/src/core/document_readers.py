"""
Document readers for various file formats including PDFs, Word docs, web pages, and text files.
Handles files of any size with memory-efficient processing.
"""

import os
import logging
from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional, Iterator
from pathlib import Path
import requests
from urllib.parse import urlparse
import tempfile

# PDF processing
import PyPDF2
import pdfplumber
from pypdf import PdfReader

# Word document processing
from docx import Document
import docx2txt

# Web scraping
import requests
from bs4 import BeautifulSoup
import trafilatura

# Text processing
import chardet
import mimetypes
from sys import platform

# Image OCR
from PIL import Image
try:
    import pytesseract
    PYTESS_AVAILABLE = True
except Exception:
    PYTESS_AVAILABLE = False
try:
    import easyocr  # optional fallback
    EASY_OCR_AVAILABLE = True
except Exception:
    EASY_OCR_AVAILABLE = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class BaseDocumentReader(ABC):
    """Abstract base class for document readers."""
    
    def __init__(self):
        self.supported_extensions = []
        # 500MB default limit; can be overridden with env var MAX_FILE_SIZE_MB
        try:
            max_mb = int(os.getenv('MAX_FILE_SIZE_MB', '500'))
        except Exception:
            max_mb = 500
        self.max_file_size = max_mb * 1024 * 1024
    
    @abstractmethod
    def read(self, source: str) -> Dict[str, Any]:
        """
        Read document and return structured data.
        
        Args:
            source: File path or URL
            
        Returns:
            Dict containing 'content', 'metadata', and 'chunks'
        """
        pass
    
    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding."""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # Read first 10KB
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8')
        except Exception:
            return 'utf-8'
    
    def _validate_file_size(self, file_path: str) -> bool:
        """Check if file size is within limits."""
        try:
            size = os.path.getsize(file_path)
            if size > self.max_file_size:
                logger.warning(f"File {file_path} exceeds size limit: {size} bytes")
                return False
            return True
        except Exception as e:
            logger.error(f"Error checking file size: {e}")
            return False


class PDFReader(BaseDocumentReader):
    """Reader for PDF documents with fallback strategies."""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pdf']
    
    def read(self, source: str) -> Dict[str, Any]:
        """Read PDF document with multiple extraction strategies."""
        if not self._validate_file_size(source):
            raise ValueError(f"PDF file too large: {source}")
        
        content = ""
        metadata = {
            'source': source,
            'type': 'pdf',
            'pages': 0,
            'extraction_method': None
        }
        
        # Try pdfplumber first (best for complex layouts)
        try:
            content, pages = self._extract_with_pdfplumber(source)
            metadata['extraction_method'] = 'pdfplumber'
            metadata['pages'] = pages
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}")
            
            # Fallback to pypdf
            try:
                content, pages = self._extract_with_pypdf(source)
                metadata['extraction_method'] = 'pypdf'
                metadata['pages'] = pages
            except Exception as e:
                logger.warning(f"pypdf failed: {e}")
                
                # Final fallback to PyPDF2
                try:
                    content, pages = self._extract_with_pypdf2(source)
                    metadata['extraction_method'] = 'PyPDF2'
                    metadata['pages'] = pages
                except Exception as e:
                    logger.error(f"All PDF extraction methods failed: {e}")
                    raise ValueError(f"Could not extract text from PDF: {source}")
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'raw_text': content.strip()
        }

    def _extract_with_pdfplumber(self, file_path: str) -> tuple[str, int]:
        """Extract text using pdfplumber (robust for complex layouts)."""
        try:
            import pdfplumber  # type: ignore
        except Exception as e:
            raise RuntimeError(f"pdfplumber not available: {e}")
        content = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                if page_text:
                    content += page_text + "\n\n"
            pages = len(pdf.pages)
        return content, pages

    def _extract_with_pypdf(self, file_path: str) -> tuple[str, int]:
        """Extract text using pypdf (lightweight and fast)."""
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception as e:
            raise RuntimeError(f"pypdf not available: {e}")
        content = ""
        reader = PdfReader(file_path)
        for page in reader.pages:
            try:
                content += (page.extract_text() or "") + "\n\n"
            except Exception:
                content += "\n\n"
        pages = len(reader.pages)
        return content, pages

    def _extract_with_pypdf2(self, file_path: str) -> tuple[str, int]:
        """Extract text using PyPDF2 as an additional fallback."""
        try:
            import PyPDF2  # type: ignore
        except Exception as e:
            raise RuntimeError(f"PyPDF2 not available: {e}")
        content = ""
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    content += (page.extract_text() or "") + "\n\n"
                except Exception:
                    content += "\n\n"
            pages = len(reader.pages)
        return content, pages


class HtmlFileReader(BaseDocumentReader):
    """Reader for local HTML files (.html/.htm)."""
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.html', '.htm']
    
    def read(self, source: str) -> Dict[str, Any]:
        if not os.path.isfile(source):
            raise ValueError(f"HTML file not found: {source}")
        if not self._validate_file_size(source):
            raise ValueError(f"HTML file too large: {source}")
        
        with open(source, 'rb') as f:
            raw = f.read()
        # Try to detect encoding
        enc = chardet.detect(raw).get('encoding') or 'utf-8'
        try:
            text = raw.decode(enc, errors='ignore')
        except Exception:
            text = raw.decode('utf-8', errors='ignore')
        
        soup = BeautifulSoup(text, 'html.parser')
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        content = soup.get_text(separator='\n')
        lines = (line.strip() for line in content.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        content = '\n'.join(chunk for chunk in chunks if chunk)
        
        metadata = {
            'source': source,
            'type': 'html',
            'encoding': enc,
            'extraction_method': 'beautifulsoup'
        }
        return {
            'content': content.strip(),
            'metadata': metadata,
            'raw_text': content.strip()
        }


class PowerPointReader(BaseDocumentReader):
    """Reader for PowerPoint files (.pptx)."""
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx']
        try:
            from pptx import Presentation  # noqa: F401
            self._pptx_available = True
        except Exception:
            self._pptx_available = False
    
    def read(self, source: str) -> Dict[str, Any]:
        if not os.path.isfile(source):
            raise ValueError(f"PPTX file not found: {source}")
        if not self._validate_file_size(source):
            raise ValueError(f"PPTX file too large: {source}")
        if not self._pptx_available:
            raise ValueError("python-pptx not installed. Install with: pip install python-pptx")
        
        from pptx import Presentation
        prs = Presentation(source)
        texts = []
        slide_idx = 0
        for slide in prs.slides:
            slide_idx += 1
            # Title and subtitle
            if slide.shapes.title and slide.shapes.title.text:
                texts.append(slide.shapes.title.text)
            # All text-containing shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                # Tables
                if shape.has_table:
                    tbl = shape.table
                    for r in tbl.rows:
                        cells_text = []
                        for c in r.cells:
                            try:
                                cells_text.append(c.text.strip())
                            except Exception:
                                pass
                        if cells_text:
                            texts.append(" | ".join(cells_text))
        content = '\n\n'.join(t.strip() for t in texts if t and t.strip())
        metadata = {
            'source': source,
            'type': 'pptx',
            'slides': len(prs.slides),
            'extraction_method': 'python-pptx'
        }
        return {'content': content, 'metadata': metadata, 'raw_text': content}


class ExcelReader(BaseDocumentReader):
    """Reader for Excel workbooks (.xlsx/.xlsm)."""
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.xlsx', '.xlsm']
        try:
            import openpyxl  # noqa: F401
            self._openpyxl_available = True
        except Exception:
            self._openpyxl_available = False
    
    def read(self, source: str) -> Dict[str, Any]:
        if not os.path.isfile(source):
            raise ValueError(f"Excel file not found: {source}")
        if not self._validate_file_size(source):
            raise ValueError(f"Excel file too large: {source}")
        if not self._openpyxl_available:
            raise ValueError("openpyxl not installed. Install with: pip install openpyxl")
        
        import openpyxl
        wb = openpyxl.load_workbook(source, read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            texts.append(f"# Sheet: {ws.title}")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_count += 1
                # To avoid huge memory, limit extremely large rows
                if row_count > 200000:  # hard safety cap
                    texts.append("[Truncated additional rows for safety]")
                    break
                values = [str(v) for v in row if v is not None]
                if values:
                    texts.append(' | '.join(values))
        content = '\n'.join(texts)
        metadata = {
            'source': source,
            'type': 'excel',
            'sheets': len(wb.worksheets),
            'extraction_method': 'openpyxl'
        }
        return {'content': content, 'metadata': metadata, 'raw_text': content}


class ImageReader(BaseDocumentReader):
    """Reader for image files using OCR."""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.png', '.jpg', '.jpeg', '.bmp', '.tif', '.tiff', '.webp', '.gif']
        # Try to locate Tesseract on Windows if not on PATH
        if PYTESS_AVAILABLE:
            try:
                # First, check explicit env var override
                tess_env = os.getenv('TESSERACT_PATH')
                if tess_env and os.path.exists(tess_env):
                    pytesseract.pytesseract.tesseract_cmd = tess_env
                elif platform.startswith('win'):
                    possible_paths = [
                        r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe",
                        r"C:\\Program Files (x86)\\Tesseract-OCR\\tesseract.exe"
                    ]
                    for p in possible_paths:
                        if os.path.exists(p):
                            pytesseract.pytesseract.tesseract_cmd = p
                            break
            except Exception:
                pass
    
    def read(self, source: str) -> Dict[str, Any]:
        if not os.path.isfile(source):
            raise ValueError(f"Image file not found: {source}")
        if not self._validate_file_size(source):
            raise ValueError(f"Image file too large: {source}")
        
        metadata = {
            'source': source,
            'type': 'image',
            'extraction_method': None,
            'ocr_engine': None
        }
        
        content = ""
        # Primary OCR: Tesseract (if available)
        used_ocr = None
        if PYTESS_AVAILABLE:
            try:
                with Image.open(source) as img:
                    # Convert to RGB for consistent OCR
                    if img.mode not in ("L", "RGB"):
                        img = img.convert("RGB")
                    content = pytesseract.image_to_string(img)
                    metadata['extraction_method'] = 'ocr'
                    metadata['ocr_engine'] = 'pytesseract'
                    used_ocr = 'pytesseract'
            except Exception as e:
                logger.warning(f"pytesseract failed: {e}")

        # Fallback to EasyOCR if needed
        if not content and EASY_OCR_AVAILABLE:
            try:
                reader = easyocr.Reader(['en'], gpu=False)
                result = reader.readtext(source, detail=0, paragraph=True)
                content = "\n".join(result)
                metadata['extraction_method'] = 'ocr'
                metadata['ocr_engine'] = 'easyocr'
                used_ocr = 'easyocr'
            except Exception as ee:
                logger.error(f"easyocr failed: {ee}")

        if not content:
            raise ValueError(
                "No OCR engine available. Install system Tesseract + pytesseract or install easyocr."
            )
        
        return {
            'content': (content or '').strip(),
            'metadata': metadata,
            'raw_text': (content or '').strip()
        }
    
    def _extract_with_pdfplumber(self, file_path: str) -> tuple[str, int]:
        """Extract text using pdfplumber (best for tables and complex layouts)."""
        content = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    content += page_text + "\n\n"
        return content, len(pdf.pages)
    
    def _extract_with_pypdf(self, file_path: str) -> tuple[str, int]:
        """Extract text using pypdf."""
        content = ""
        reader = PdfReader(file_path)
        for page in reader.pages:
            content += page.extract_text() + "\n\n"
        return content, len(reader.pages)
    
    def _extract_with_pypdf2(self, file_path: str) -> tuple[str, int]:
        """Extract text using PyPDF2 as final fallback."""
        content = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                content += page.extract_text() + "\n\n"
        return content, len(pdf_reader.pages)


class WordDocumentReader(BaseDocumentReader):
    """Reader for Word documents (.docx and .doc)."""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.docx', '.doc']
    
    def read(self, source: str) -> Dict[str, Any]:
        """Read Word document."""
        if not self._validate_file_size(source):
            raise ValueError(f"Word document too large: {source}")
        
        content = ""
        metadata = {
            'source': source,
            'type': 'word',
            'paragraphs': 0,
            'extraction_method': None
        }
        
        file_ext = Path(source).suffix.lower()
        
        if file_ext == '.docx':
            try:
                # Try python-docx first
                doc = Document(source)
                paragraphs = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text)
                
                content = '\n\n'.join(paragraphs)
                metadata['paragraphs'] = len(paragraphs)
                metadata['extraction_method'] = 'python-docx'
                
            except Exception as e:
                logger.warning(f"python-docx failed: {e}")
                # Fallback to docx2txt
                try:
                    content = docx2txt.process(source)
                    metadata['extraction_method'] = 'docx2txt'
                except Exception as e:
                    logger.error(f"docx2txt failed: {e}")
                    raise ValueError(f"Could not extract text from Word document: {source}")
        
        elif file_ext == '.doc':
            # For .doc files, we need additional tools like antiword or textract
            # For now, we'll raise an informative error
            raise ValueError(f".doc files require additional tools. Please convert to .docx format: {source}")
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'raw_text': content.strip()
        }


class WebPageReader(BaseDocumentReader):
    """Reader for web pages and URLs."""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = []  # URLs don't have extensions
        self.timeout = 30
        self.headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
    
    def read(self, source: str) -> Dict[str, Any]:
        """Read web page content."""
        if not self._is_valid_url(source):
            raise ValueError(f"Invalid URL: {source}")
        
        content = ""
        metadata = {
            'source': source,
            'type': 'webpage',
            'title': '',
            'extraction_method': None
        }
        
        # Try trafilatura first (best for article extraction)
        try:
            downloaded = trafilatura.fetch_url(source)
            if downloaded:
                content = trafilatura.extract(downloaded, include_comments=False, include_tables=True)
                if content:
                    metadata['extraction_method'] = 'trafilatura'
                    # Try to get title
                    soup = BeautifulSoup(downloaded, 'html.parser')
                    title_tag = soup.find('title')
                    if title_tag:
                        metadata['title'] = title_tag.get_text().strip()
        except Exception as e:
            logger.warning(f"trafilatura failed: {e}")
        
        # Fallback to BeautifulSoup
        if not content:
            try:
                response = requests.get(source, headers=self.headers, timeout=self.timeout)
                response.raise_for_status()
                
                soup = BeautifulSoup(response.content, 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                
                # Get title
                title_tag = soup.find('title')
                if title_tag:
                    metadata['title'] = title_tag.get_text().strip()
                
                # Extract text
                content = soup.get_text()
                
                # Clean up text
                lines = (line.strip() for line in content.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                content = '\n'.join(chunk for chunk in chunks if chunk)
                
                metadata['extraction_method'] = 'beautifulsoup'
                
            except Exception as e:
                logger.error(f"Web scraping failed: {e}")
                raise ValueError(f"Could not extract content from URL: {source}")
        
        if not content:
            raise ValueError(f"No content extracted from URL: {source}")
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'raw_text': content.strip()
        }
    
    def _is_valid_url(self, url: str) -> bool:
        """Check if URL is valid."""
        try:
            result = urlparse(url)
            return all([result.scheme, result.netloc])
        except Exception:
            return False


class TextFileReader(BaseDocumentReader):
    """Reader for plain text files."""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.txt', '.md', '.csv', '.json', '.xml', '.log']
    
    def read(self, source: str) -> Dict[str, Any]:
        """Read text file with encoding detection."""
        if not self._validate_file_size(source):
            raise ValueError(f"Text file too large: {source}")
        
        encoding = self._detect_encoding(source)
        
        try:
            with open(source, 'r', encoding=encoding) as file:
                content = file.read()
        except UnicodeDecodeError:
            # Try with different encodings
            for enc in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(source, 'r', encoding=enc) as file:
                        content = file.read()
                    encoding = enc
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError(f"Could not decode text file: {source}")
        
        metadata = {
            'source': source,
            'type': 'text',
            'encoding': encoding,
            'size_bytes': os.path.getsize(source),
            'lines': content.count('\n') + 1
        }
        
        return {
            'content': content.strip(),
            'metadata': metadata,
            'raw_text': content.strip()
        }


class DocumentReaderFactory:
    """Factory class to get appropriate document reader."""
    
    def __init__(self):
        # Map logical types to reader classes (instantiate per request to avoid state carryover)
        self.reader_classes = {
            'pdf': PDFReader,
            'word': WordDocumentReader,
            'webpage': WebPageReader,
            'text': TextFileReader,
            'image': ImageReader,
            'html': HtmlFileReader,
            'pptx': PowerPointReader,
            'excel': ExcelReader,
        }
    
    def get_reader(self, source: str) -> BaseDocumentReader:
        """Get appropriate reader based on source type."""
        
        # Check if it's a URL
        if self._is_url(source):
            return self.reader_classes['webpage']()
        
        # Check file extension
        if os.path.isfile(source):
            file_ext = Path(source).suffix.lower()
            
            if file_ext == '.pdf':
                return self.reader_classes['pdf']()
            elif file_ext in ['.docx', '.doc']:
                return self.reader_classes['word']()
            elif file_ext in ['.txt', '.md', '.csv', '.json', '.xml', '.log']:
                return self.reader_classes['text']()
            elif file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tif', '.tiff', '.webp', '.gif']:
                return self.reader_classes['image']()
            elif file_ext in ['.html', '.htm']:
                return self.reader_classes['html']()
            elif file_ext in ['.pptx']:
                return self.reader_classes['pptx']()
            elif file_ext in ['.xlsx', '.xlsm']:
                return self.reader_classes['excel']()
            elif file_ext in ['.xls']:
                raise ValueError("Legacy .xls not supported without xlrd; please convert to .xlsx or install xlrd==1.2.0")
            else:
                # Try to determine by MIME type
                mime_type, _ = mimetypes.guess_type(source)
                if mime_type:
                    if 'pdf' in mime_type:
                        return self.reader_classes['pdf']()
                    elif 'word' in mime_type or 'document' in mime_type:
                        return self.reader_classes['word']()
                    elif 'text' in mime_type:
                        return self.reader_classes['text']()
                    elif mime_type.startswith('image/'):
                        return self.reader_classes['image']()
                    elif mime_type in ['text/html', 'application/xhtml+xml']:
                        return self.reader_classes['html']()
                    elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                        return self.reader_classes['pptx']()
                    elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                        return self.reader_classes['excel']()
                
                # Default to text reader for unknown extensions
                logger.warning(f"Unknown file type, defaulting to text reader: {source}")
                return self.reader_classes['text']()
        
        raise ValueError(f"Unsupported source type: {source}")
    
    def _is_url(self, source: str) -> bool:
        """Check if source is a URL."""
        try:
            result = urlparse(source)
            return all([result.scheme, result.netloc])
        except Exception:
            return False
    
    def read_document(self, source: str) -> Dict[str, Any]:
        """Read document using appropriate reader."""
        reader = self.get_reader(source)
        return reader.read(source)
